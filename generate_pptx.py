"""
VisionDX Mega — REGENERATE PPTX with improved architecture slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── COLORS ──────────────────────────────────────────────────────────────────
DARK_BG       = RGBColor(0x02, 0x08, 0x10)
CARD_BG       = RGBColor(0x0B, 0x13, 0x29)
CARD2_BG      = RGBColor(0x0D, 0x1A, 0x35)
ACCENT_GREEN  = RGBColor(0x00, 0xD4, 0xAA)
ACCENT_BLUE   = RGBColor(0x00, 0x9F, 0xFF)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_ROSE   = RGBColor(0xF4, 0x3F, 0x5E)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_70      = RGBColor(0xB3, 0xBF, 0xCC)
TEAL          = RGBColor(0x14, 0xB8, 0xA6)
EMERALD       = RGBColor(0x10, 0xB9, 0x81)
CYAN          = RGBColor(0x06, 0xB6, 0xD4)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def add_rrect(slide, x, y, w, h, fill, border=None, border_w=1.0, corner=0.06):
    s = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
    else:
        s.line.fill.background()
    s.adjustments[0] = corner
    return s

def add_text(slide, text, x, y, w, h, size=14, color=WHITE, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb

def add_multiline(slide, lines, x, y, w, h, size=12, color=WHITE_70, bold=False, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(2)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font
    return tb

def add_pill(slide, text, x, y, w, h, bg, fg=WHITE, size=10):
    s = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = bg
    s.line.color.rgb = fg; s.line.width = Pt(0.5)
    s.adjustments[0] = 0.5
    tf = s.text_frame; tf.text = text
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(size); p.runs[0].font.color.rgb = fg
    p.runs[0].font.bold = True; p.runs[0].font.name = "Calibri"
    return s

def add_bg(slide):
    bg = add_rect(slide, 0, 0, 13.333, 7.5, DARK_BG)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

def top_bar(slide, col=ACCENT_GREEN):
    b = add_rect(slide, 0, 0, 13.333, 0.055, col)
    return b

def bottom_bar(slide):
    add_rect(slide, 0, 7.32, 13.333, 0.18, CARD_BG)
    add_text(slide, "VisionDX Mega Platform  •  Pakistan Ka Digital Doctor  •  Hackathon 2025",
             0, 7.33, 13.333, 0.18, size=8.5, color=WHITE_70, align=PP_ALIGN.CENTER)

def slide_header(slide, badge, title, sub=None, badge_col=ACCENT_GREEN, accent=ACCENT_GREEN):
    top_bar(slide, accent)
    bottom_bar(slide)
    add_pill(slide, badge, 0.35, 0.18, 2.5, 0.28, CARD_BG, badge_col, 8.5)
    add_text(slide, title, 0.35, 0.52, 12.7, 0.72, size=32, color=WHITE, bold=True)
    line = add_rect(slide, 0.35, 1.26, 2.8, 0.04, accent)
    if sub:
        add_text(slide, sub, 0.35, 1.32, 12.7, 0.35, size=12.5, color=WHITE_70)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
top_bar(sl, ACCENT_GREEN)
bottom_bar(sl)

# Hero circle
c = sl.shapes.add_shape(9, Inches(0.45), Inches(1.1), Inches(2.0), Inches(2.0))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x00,0x1E,0x18)
c.line.color.rgb = ACCENT_GREEN; c.line.width = Pt(2.5)
add_text(sl, "🏥", 0.65, 1.35, 1.6, 1.2, size=48, align=PP_ALIGN.CENTER)

add_pill(sl, "HACKATHON 2025  •  AI FOR HEALTHCARE  •  PAKISTAN", 2.65, 1.15, 8.5, 0.34,
         RGBColor(0x0A,0x22,0x1A), ACCENT_GREEN, 9.5)

add_text(sl, "VisionDX Mega", 2.65, 1.6, 10.5, 1.0, size=56, color=WHITE, bold=True)
add_text(sl, "Platform", 2.65, 2.5, 10.5, 0.9, size=56, color=ACCENT_GREEN, bold=True)
add_text(sl, "Pakistan Ka Digital Doctor  •  پاکستان کا ڈیجیٹل ڈاکٹر",
         2.65, 3.48, 10.5, 0.4, size=15.5, color=ACCENT_BLUE)

div = add_rect(sl, 2.65, 3.97, 7.5, 0.04, ACCENT_GREEN)
add_text(sl, "12 AI-Powered Diagnostics  •  100% Private  •  Offline-Ready  •  7 Languages",
         2.65, 4.06, 10.5, 0.4, size=12, color=WHITE_70)

for i,(num,lbl,col) in enumerate([("12","AI Modules",ACCENT_GREEN),("7+","Languages",ACCENT_BLUE),("0","Data Leaks",ACCENT_PURPLE)]):
    bx = 2.65 + i*3.35
    b = add_rrect(sl, bx, 4.62, 3.1, 1.18, CARD_BG, col, 1.2, 0.08)
    add_text(sl, num, bx+0.1, 4.65, 2.9, 0.66, size=32, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, bx+0.1, 5.26, 2.9, 0.35, size=10.5, color=WHITE_70, align=PP_ALIGN.CENTER)

add_text(sl, "React 18 · Vite · TailwindCSS · Groq LPU · Whisper V3 · OpenCV · MediaPipe",
         0.35, 6.78, 12.7, 0.38, size=9.5, color=WHITE_70, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
slide_header(sl, "01  THE PROBLEM", "Pakistan's Healthcare Crisis",
             "240 million people. Chronic under-investment. Growing AI opportunity.",
             ACCENT_ROSE, ACCENT_ROSE)

stats = [
    ("73%","Population in rural areas with NO specialist access",ACCENT_ROSE),
    ("1 : 1,600","Doctor-to-patient ratio (WHO recommends 1:1,000)",ACCENT_AMBER),
    ("68%","Cannot afford private clinic consultations",ACCENT_PURPLE),
    ("42%","Low health literacy — cannot interpret lab reports",ACCENT_BLUE),
]
for i,(num,desc,col) in enumerate(stats):
    row,c = i//2, i%2
    px,py = 0.35 + c*6.55, 1.73 + row*2.62
    b = add_rrect(sl, px, py, 6.35, 2.38, CARD_BG, col, 2.0, 0.06)
    # Left accent strip
    strip = add_rect(sl, px, py, 0.08, 2.38, col)
    add_text(sl, num, px+0.22, py+0.15, 6.0, 0.85, size=40, color=col, bold=True)
    add_text(sl, desc, px+0.22, py+0.98, 5.9, 1.22, size=13, color=WHITE_70)

pains = ["Language barriers → misdiagnosis risk (Urdu, Pashto, Sindhi, Punjabi patients unserved)",
         "No digital wound tracking for 100,000+ Lady Health Workers in field settings",
         "Pakistan's 10M+ Deaf community completely excluded from medical communication systems"]
add_text(sl, "Root Causes:", 0.35, 6.1, 3, 0.3, size=11, color=ACCENT_ROSE, bold=True)
add_multiline(sl, ["  ✕  "+p for p in pains], 0.35, 6.4, 12.8, 0.9, size=10, color=WHITE_70)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
slide_header(sl, "02  OUR SOLUTION", "VisionDX Mega — One Platform, 12 AI Doctors",
             "Unified client-side AI diagnostics for every Pakistani patient, online or offline",
             ACCENT_GREEN, ACCENT_GREEN)

# Hub
hub = sl.shapes.add_shape(9, Inches(5.4), Inches(2.4), Inches(2.55), Inches(2.55))
hub.fill.solid(); hub.fill.fore_color.rgb = RGBColor(0x00,0x1E,0x18)
hub.line.color.rgb = ACCENT_GREEN; hub.line.width = Pt(2.5)
add_text(sl, "🏥", 5.4, 2.45, 2.55, 1.0, size=30, align=PP_ALIGN.CENTER)
add_text(sl, "VisionDX\nMega", 5.4, 3.35, 2.55, 0.9, size=15, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

modules = [
    ("👁️ Eye AI",0.3,2.3,ACCENT_BLUE), ("🔴 Skin AI",0.3,3.55,ACCENT_ROSE),
    ("🩹 Wound CV",0.3,4.8,ACCENT_GREEN), ("🎤 Cough",2.0,1.55,ACCENT_BLUE),
    ("😴 Sleep AI",4.1,1.55,ACCENT_PURPLE), ("📄 Lab AI",7.2,1.55,CYAN),
    ("📅 Routine",9.2,1.55,ACCENT_PURPLE), ("💊 Medicine",9.8,2.3,ACCENT_AMBER),
    ("💇 Hair AI",9.8,3.55,EMERALD), ("💊 Rx AI",9.8,4.8,ACCENT_AMBER),
    ("🩺 VoiceDoc",2.1,5.7,EMERALD), ("🖐️ PSL Sign",5.3,5.7,TEAL),
]
for (lbl,mx,my,mc) in modules:
    mb = add_rrect(sl, mx, my, 2.05, 0.6, CARD_BG, mc, 0.9, 0.5)
    add_text(sl, lbl, mx+0.08, my+0.1, 1.9, 0.44, size=10, color=mc, bold=True, align=PP_ALIGN.CENTER)

diffs = ["  ✅  100% Client-side — Zero backend patient data", "  ✅  Offline Service Worker for rural areas",
         "  ✅  Urdu/Punjabi/Pashto/Sindhi voice responses", "  ✅  WhatsApp Bot for zero-smartphone patients"]
add_text(sl, "Key Differentiators:", 0.35, 6.1, 5, 0.3, size=11, color=ACCENT_GREEN, bold=True)
add_multiline(sl, diffs, 0.35, 6.4, 12.8, 0.9, size=10, color=WHITE_70)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE DIAGRAM (REDESIGNED)
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
top_bar(sl, ACCENT_BLUE)
bottom_bar(sl)

# Title
add_pill(sl, "03  SYSTEM ARCHITECTURE", 0.35, 0.18, 3.5, 0.28, CARD_BG, ACCENT_BLUE, 8.5)
add_text(sl, "Multi-Layer AI Architecture", 0.35, 0.52, 12.7, 0.62, size=30, color=WHITE, bold=True)
add_rect(sl, 0.35, 1.16, 2.8, 0.04, ACCENT_BLUE)
add_text(sl, "Patient Input  →  AI Processing  →  Backend Services  →  Private Output  →  Zero Data Leaks",
         0.35, 1.22, 12.7, 0.32, size=11, color=WHITE_70)

# ── LAYER LABELS ──
layer_data = [
    (0.35, 1.62, 12.65, "① INPUT LAYER — 7 Entry Points", ACCENT_BLUE, RGBColor(0x0A,0x14,0x28)),
    (0.35, 2.98, 12.65, "② AI ENGINE — Groq LPU · Gemini · MediaPipe · ElevenLabs TTS", ACCENT_GREEN, RGBColor(0x02,0x0E,0x08)),
    (0.35, 4.35, 12.65, "③ BACKEND SERVICES — FastAPI CV  +  Node.js WhatsApp Bot", ACCENT_AMBER, RGBColor(0x16,0x10,0x02)),
    (0.35, 5.62, 12.65, "④ OUTPUT — 100% Private · localStorage Only · Zero Server Storage", ACCENT_PURPLE, RGBColor(0x12,0x08,0x1E)),
]

for (lx,ly,lw,lt,lc,lb) in layer_data:
    lh = 1.25 if ly < 5 else 1.15
    box = add_rrect(sl, lx, ly, lw, lh, lb, lc, 0.8, 0.04)
    # Layer title pill
    add_pill(sl, lt, lx+0.1, ly+0.07, lw-0.2, 0.24, DARK_BG, lc, 8)

# ── INPUT items ──
inputs = [
    ("📷 Photo",ACCENT_BLUE), ("🎙️ Voice",ACCENT_GREEN), ("📸 Webcam",TEAL),
    ("✍️ Manual",ACCENT_AMBER), ("💬 Text",ACCENT_PURPLE), ("🖐️ Gesture",ACCENT_ROSE), ("💬 WhatsApp",EMERALD),
]
for i,(lbl,ic) in enumerate(inputs):
    ix = 0.45 + i*1.8
    ib = add_rrect(sl, ix, 1.93, 1.7, 0.56, CARD_BG, ic, 0.7, 0.5)
    add_text(sl, lbl, ix+0.05, 1.99, 1.6, 0.44, size=9.5, color=ic, bold=True, align=PP_ALIGN.CENTER)

# ── AI ENGINE items ──
ai_items = [
    ("⚡ Groq\nLlama 3.3 70B\nMedical Reasoning",ACCENT_GREEN,3.08),
    ("🎤 Groq\nWhisper V3\nVoice-to-Text",ACCENT_GREEN,3.08),
    ("✨ Google\nGemini 1.5\nVision Analysis",ACCENT_BLUE,3.08),
    ("🖐️ MediaPipe\nHands CDN\nSkeletal Track",TEAL,3.08),
]
for i,(lbl,ic,lh) in enumerate(ai_items):
    ax = 0.45 + i*3.2
    ab = add_rrect(sl, ax, 3.28, 3.0, 0.88, CARD2_BG, ic, 1.2, 0.07)
    add_text(sl, lbl, ax+0.12, 3.3, 2.8, 0.86, size=9.5, color=ic, bold=True, align=PP_ALIGN.CENTER)

# ── BACKEND items ──
backends = [
    ("🔬 FastAPI + OpenCV\nWound Segmentation Backend\nlocalhost:8000", ACCENT_AMBER, 4.2),
    ("💬 Node.js Express\nWhatsApp Triage Bot\nlocalhost:3001", EMERALD, 4.2),
    ("🔊 ElevenLabs TTS\nMultilingual Voice\nNeural Speech Output", ACCENT_PURPLE, 4.2),
]
for i,(lbl,ic,bh) in enumerate(backends):
    bx = 0.45 + i*4.3
    bb = add_rrect(sl, bx, 4.65, 4.1, 0.82, CARD2_BG, ic, 1.2, 0.07)
    add_text(sl, lbl, bx+0.12, 4.67, 3.9, 0.78, size=9.5, color=ic, bold=True, align=PP_ALIGN.CENTER)

# ── OUTPUT items ──
out_items = [
    ("📋 AI Medical Report",ACCENT_BLUE), ("🔊 TTS Voice Response",ACCENT_PURPLE),
    ("💬 WhatsApp Alert",EMERALD), ("🖨️ Printable PDF",ACCENT_AMBER), ("🔒 localStorage",ACCENT_GREEN),
]
for i,(lbl,ic) in enumerate(out_items):
    ox = 0.45 + i*2.56
    ob = add_rrect(sl, ox, 5.9, 2.44, 0.55, CARD_BG, ic, 0.8, 0.5)
    add_text(sl, lbl, ox+0.06, 5.95, 2.32, 0.44, size=9, color=ic, bold=True, align=PP_ALIGN.CENTER)

# ── DOWN ARROWS between layers ──
for ay in [2.52, 3.88, 5.5]:
    for ax_frac in [0.27, 0.5, 0.73]:
        arr_x = ax_frac * 13.333
        arr_box = add_rect(sl, arr_x-0.02, ay, 0.04, 0.2, ACCENT_GREEN)

# Arrow text labels
add_text(sl, "▼  ▼  ▼", 5.5, 2.52, 2, 0.25, size=9, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_text(sl, "▼  ▼  ▼", 5.5, 3.88, 2, 0.25, size=9, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)
add_text(sl, "▼  ▼  ▼", 5.5, 5.48, 2, 0.25, size=9, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — WOUND TRACKER DEEP DIVE
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
slide_header(sl, "04  DEEP DIVE", "Wound Healing Tracker — Dual Engine System",
             "OpenCV Computer Vision  +  Groq LLaMA Clinical Reasoning",
             ACCENT_GREEN, ACCENT_GREEN)

lbox = add_rrect(sl, 0.35, 1.65, 5.9, 4.75, CARD_BG, ACCENT_GREEN, 1.5, 0.05)
add_text(sl, "🔬  Computer Vision Pipeline (FastAPI + OpenCV)", 0.5, 1.72, 5.7, 0.38, size=13, color=ACCENT_GREEN, bold=True)

cv_steps = [
    ("1","Image Decode → HSV Color Space Conversion"),
    ("2","3-Mask Wound Tissue Isolation\n(Red + Deep-Red + Necrotic HSV ranges)"),
    ("3","Morphological Cleanup (Elliptic kernel close/open)"),
    ("4","Contour Detection → Largest Contour Selection"),
    ("5","Ellipse Fitting → Major/Minor Axis Measurement"),
    ("6","px → mm Calibration (100mm reference field)"),
    ("7","Healing Stage Classification by Color Analysis"),
    ("8","Segmentation Overlay → Base64 PNG Export"),
]
for i,(num,step) in enumerate(cv_steps):
    py = 2.15 + i*0.49
    pill_c = sl.shapes.add_shape(9, Inches(0.5), Inches(py+0.03), Inches(0.28), Inches(0.28))
    pill_c.fill.solid(); pill_c.fill.fore_color.rgb = ACCENT_GREEN; pill_c.line.fill.background()
    add_text(sl, num, 0.5, py+0.01, 0.28, 0.3, size=8, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, step, 0.85, py, 5.25, 0.46, size=9, color=WHITE_70)

rbox = add_rrect(sl, 6.5, 1.65, 6.5, 4.75, CARD_BG, ACCENT_BLUE, 1.5, 0.05)
add_text(sl, "📊  Output Metrics & Clinical Intelligence", 6.65, 1.72, 6.2, 0.38, size=13, color=ACCENT_BLUE, bold=True)

metrics = [
    ("Wound Area","mm² isolated tissue field",ACCENT_GREEN),
    ("Length × Width","Calibrated major/minor axis in mm",ACCENT_BLUE),
    ("Estimated Depth","Infiltration index from ellipse dims",ACCENT_AMBER),
    ("Healing Stage","Inflammatory / Proliferative / Remodeling",ACCENT_PURPLE),
    ("CV Confidence","Detection confidence score (0.0 → 1.0)",TEAL),
    ("Seg Overlay","OpenCV contour annotation PNG",ACCENT_ROSE),
]
for i,(m,d,mc) in enumerate(metrics):
    my = 2.15 + i*0.68
    mb = add_rrect(sl, 6.65, my, 6.2, 0.6, DARK_BG, mc, 0.8, 0.05)
    add_text(sl, m, 6.82, my+0.06, 2.2, 0.48, size=11, color=mc, bold=True)
    add_text(sl, d, 9.12, my+0.1, 3.6, 0.42, size=9.5, color=WHITE_70)

add_text(sl, "LLM CLINICAL LAYER: Combines CV metrics + pain level + wound location + day → Evidence-based dressing recommendations + NPWT candidacy evaluation",
         0.35, 6.44, 12.8, 0.4, size=10.5, color=ACCENT_AMBER, bold=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — VOICEDOC + PSL
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
slide_header(sl, "05  ACCESSIBILITY AI", "VoiceDoc + PSL Sign Language Translator",
             "Bridging the gap for rural patients and Pakistan's 10M+ Deaf community",
             EMERALD, EMERALD)

vbox = add_rrect(sl, 0.35, 1.65, 6.2, 4.75, CARD_BG, EMERALD, 1.5, 0.05)
add_text(sl, "🩺  VoiceDoc — Multilingual Rural AI Triage", 0.5, 1.73, 5.95, 0.38, size=13.5, color=EMERALD, bold=True)
vd = [("🎙️","Speak symptoms in Urdu/Punjabi/Pashto/Sindhi/Hindi"),
      ("⚡","Groq Whisper V3 → sub-200ms multilingual transcription"),
      ("🧠","Llama 3.3 70B → Triage + urgency classification"),
      ("🔊","ElevenLabs / Google TTS → native language response"),
      ("💬","WhatsApp Mode → structured referral → doctor alert"),]
for i,(icon,step) in enumerate(vd):
    py = 2.2 + i*0.75
    add_text(sl, icon, 0.5, py+0.05, 0.5, 0.55, size=18, align=PP_ALIGN.CENTER)
    add_text(sl, step, 1.1, py+0.08, 5.3, 0.55, size=10.5, color=WHITE_70)

langs = ["Urdu اردو","Roman Urdu","Punjabi پنجابی","Pashto پښتو","Sindhi سنڌي","Hindi हिन्दी","English"]
add_text(sl, "Languages:", 0.5, 5.95, 2.5, 0.28, size=10, color=WHITE, bold=True)
for i,lang in enumerate(langs):
    lx = 0.5 + (i%4)*1.5
    ly = 6.25 + (i//4)*0.35
    add_pill(sl, lang, lx, ly, 1.4, 0.28, RGBColor(0x00,0x1E,0x18), EMERALD, 8)

pbox = add_rrect(sl, 6.75, 1.65, 6.25, 4.75, CARD_BG, TEAL, 1.5, 0.05)
add_text(sl, "🖐️  PSL Sign Language Translator", 6.9, 1.73, 5.95, 0.38, size=13.5, color=TEAL, bold=True)
psl_f = [("MediaPipe Hands","21-landmark skeletal tracking at 30FPS"),
         ("PSL Dictionary","60+ signs: Greetings, Household, Sentences"),
         ("Urdu Alphabets","Complete Urdu alphabet (Alef → Ye) mapping"),
         ("Dual-Hand Signs","Two-hand compound gesture support"),
         ("AI Sentence AI","Gemini/Groq compiles signs → full sentence"),
         ("Learning Center","Gamified teaching with live camera validation"),
         ("Data Collector","Export landmark JSON for model training"),]
for i,(feat,desc) in enumerate(psl_f):
    fy = 2.18 + i*0.595
    fb = add_rrect(sl, 6.9, fy, 6.0, 0.52, DARK_BG, TEAL, 0.6, 0.05)
    add_text(sl, feat, 7.06, fy+0.06, 2.0, 0.42, size=10, color=TEAL, bold=True)
    add_text(sl, desc, 9.18, fy+0.09, 3.6, 0.38, size=9, color=WHITE_70)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — ALL 12 MODULES
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
slide_header(sl, "06  ALL 12 MODULES", "Comprehensive Diagnostic Module Suite",
             "Every tool designed to replace an expensive specialist visit",
             ACCENT_BLUE, ACCENT_BLUE)

all_mods = [
    ("👁️","Eye Disease\nPredictor","Conjunctivitis, Cataracts,\nGlaucoma, Retinopathy",ACCENT_BLUE),
    ("🔴","Skin Rash\nAnalyzer","Melanoma, BCC, Eczema,\nPsoriasis, Tinea",ACCENT_ROSE),
    ("🩹","Wound\nHealer CV","OpenCV segmentation,\nmm² tracking, staging",ACCENT_GREEN),
    ("🎤","Cough\nSound AI","Dry/wet/pertussis,\nWhisper V3 audio",ACCENT_BLUE),
    ("😴","Sleep\nQuality AI","Apnea risk, STOP-BANG,\nEpworth scoring",ACCENT_PURPLE),
    ("💊","Medicine\nVerifier","Counterfeit detection,\nOCR, dosage guide",ACCENT_AMBER),
    ("📄","Lab Report\nAnalyzer","CBC, Lipid, HbA1c,\nLFT, KFT + print",CYAN),
    ("💇","Hair Disease\nAI","Alopecia, trichology,\nscalp conditions",EMERALD),
    ("📅","Daily Routine\nAnalyzer","Sleep, screen, water,\nstress optimizer",ACCENT_PURPLE),
    ("🩺","VoiceDoc\nTriage","7-language AI doctor,\nWhatsApp forward",EMERALD),
    ("🖐️","PSL Sign\nTranslator","MediaPipe + 60+ PSL\nsigns + Urdu abc",TEAL),
    ("💊","Suggest\nMedicine AI","Rx pad, OTC, dosages,\nblackbox warnings",ACCENT_AMBER),
]
for i,(icon,name,desc,col) in enumerate(all_mods):
    row,c = i//4, i%4
    mx = 0.3 + c*3.26
    my = 1.65 + row*1.85
    mc = add_rrect(sl, mx, my, 3.1, 1.7, CARD_BG, col, 1.0, 0.06)
    add_text(sl, icon, mx+0.1, my+0.06, 0.62, 0.62, size=24, align=PP_ALIGN.CENTER)
    add_text(sl, name, mx+0.76, my+0.06, 2.22, 0.6, size=11, color=col, bold=True)
    add_text(sl, desc, mx+0.12, my+0.78, 2.9, 0.78, size=8.5, color=WHITE_70)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — CHW + WHATSAPP
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
slide_header(sl, "07  RURAL IMPACT", "CHW Platform & WhatsApp Triage Bot",
             "Empowering Lady Health Workers and zero-smartphone communities",
             ACCENT_AMBER, ACCENT_AMBER)

cbox = add_rrect(sl, 0.35, 1.65, 6.2, 4.75, CARD_BG, ACCENT_AMBER, 1.5, 0.05)
add_text(sl, "🧳  Community Health Worker (CHW) Platform", 0.5, 1.72, 5.9, 0.38, size=13, color=ACCENT_AMBER, bold=True)
chw_f = [
    "Dedicated CHW login flow separate from patient portal",
    "Risk-stratified patient queue:\n   Critical  •  Warning  •  Normal",
    "Auto-saves ALL 12 module results via useSaveToCHW() hook",
    "Generates printable A4 patient summary for hospital handoffs",
    "Field-optimized offline-ready interface for rural clinics",
    "Serves 100,000+ Lady Health Workers nationwide",
]
for i,feat in enumerate(chw_f):
    fy = 2.22 + i*0.7
    add_text(sl, "▸", 0.5, fy+0.06, 0.3, 0.55, size=12, color=ACCENT_AMBER, bold=True)
    add_text(sl, feat, 0.82, fy, 5.55, 0.68, size=10, color=WHITE_70)

wbox = add_rrect(sl, 6.75, 1.65, 6.25, 4.75, CARD_BG, EMERALD, 1.5, 0.05)
add_text(sl, "💬  WhatsApp Triage Bot", 6.9, 1.72, 6.0, 0.38, size=13, color=EMERALD, bold=True)
wa = [
    ("Patient texts symptoms\n(any language, zero app needed", EMERALD),
    ("VoiceDoc AI → triage analysis\n+ urgency classification", ACCENT_GREEN),
    ("Structured referral via\nMeta Business API dispatch", ACCENT_BLUE),
    ("Doctor receives formatted\nWhatsApp alert + patient profile", ACCENT_AMBER),
]
for i,(step,sc) in enumerate(wa):
    sy = 2.2 + i*0.95
    sb = add_rrect(sl, 6.9, sy, 6.0, 0.78, DARK_BG, sc, 1.0, 0.06)
    add_text(sl, f"Step {i+1}", 7.05, sy+0.06, 1.1, 0.28, size=9, color=sc, bold=True)
    add_text(sl, step, 7.05, sy+0.3, 5.7, 0.44, size=10, color=WHITE_70)

prev = add_rrect(sl, 6.9, 6.0, 6.0, 0.86, RGBColor(0x00,0x1E,0x18), EMERALD, 1, 0.05)
add_text(sl, "VISIONDX TRIAGE ALERT  |  Patient: Abdullah (18M)  |  Rahim Yar Khan\nPriority: SEE_DOCTOR  •  Groq AI Recommendation: Seek immediate care",
         7.02, 6.04, 5.8, 0.82, size=8.5, color=EMERALD)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECH STACK + PRIVACY
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
slide_header(sl, "08  TECH STACK", "Technology Stack & Privacy Architecture",
             "Best-in-class open tools assembled for maximum impact at zero patient cost",
             ACCENT_PURPLE, ACCENT_PURPLE)

cats = [
    ("Frontend", [("React 18.3","Core UI",ACCENT_BLUE),("Vite 5.2","Build+HMR",ACCENT_PURPLE),
                  ("TailwindCSS 3.4","Styling",ACCENT_BLUE),("React Router 6","SPA Routing",TEAL)], ACCENT_BLUE),
    ("AI & APIs", [("Groq Llama 3.3","Medical Text",ACCENT_GREEN),("Groq Whisper V3","Voice STT",ACCENT_GREEN),
                   ("Gemini 1.5","Vision AI",ACCENT_BLUE),("ElevenLabs TTS","Neural Voice",ACCENT_PURPLE)], ACCENT_GREEN),
    ("Backends", [("FastAPI 0.110","Wound CV API",ACCENT_AMBER),("OpenCV 4.9","Segmentation",ACCENT_AMBER),
                  ("Node.js Express","WhatsApp Bot",EMERALD),("MediaPipe CDN","Hand Track",TEAL)], ACCENT_AMBER),
]
for ci,(cat,items,cc) in enumerate(cats):
    cx = 0.3 + ci*4.42
    cbox = add_rrect(sl, cx, 1.65, 4.22, 4.0, CARD_BG, cc, 1.5, 0.05)
    add_text(sl, cat, cx+0.15, 1.73, 3.95, 0.36, size=13, color=cc, bold=True)
    for ii,(tech,role,tc) in enumerate(items):
        iy = 2.2 + ii*0.83
        ib = add_rrect(sl, cx+0.15, iy, 3.95, 0.68, DARK_BG, tc, 0.7, 0.06)
        add_text(sl, tech, cx+0.3, iy+0.06, 2.5, 0.32, size=11, color=tc, bold=True)
        add_text(sl, role, cx+0.3, iy+0.36, 3.5, 0.28, size=9, color=WHITE_70)

pbox = add_rrect(sl, 13.3-4.25, 1.65, 4.25, 4.0, CARD_BG, ACCENT_GREEN, 2.0, 0.05)
add_text(sl, "Privacy", 13.3-4.1, 1.73, 3.95, 0.36, size=13, color=ACCENT_GREEN, bold=True)
priv = [("Zero server-side","patient data storage"),("All records in","browser localStorage"),
        ("API keys","stored client-side only"),("No telemetry","or analytics tracking"),
        ("CORS-protected","backends (localhost)"),]
for i,(p1,p2) in enumerate(priv):
    py = 2.2 + i*0.73
    add_text(sl, "  "+p1, 13.3-4.08, py+0.02, 3.85, 0.3, size=10, color=ACCENT_GREEN, bold=True)
    add_text(sl, "  "+p2, 13.3-4.08, py+0.3, 3.85, 0.3, size=9, color=WHITE_70)

add_text(sl, "Medical Disclaimer: VisionDX is an AI assistance tool. Always consult a licensed physician for clinical decisions.",
         0.35, 5.85, 12.8, 0.38, size=10, color=ACCENT_AMBER, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — IMPACT + CALL TO ACTION
# ═══════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
add_bg(sl)
top_bar(sl, ACCENT_GREEN)
bottom_bar(sl)

add_pill(sl, "HACKATHON 2025  •  AI FOR HEALTHCARE  •  REAL IMPACT", 1.5, 0.2, 10.3, 0.34,
         CARD_BG, ACCENT_GREEN, 9)
add_text(sl, "Real Impact.", 0.4, 0.68, 12.6, 0.95, size=50, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, "Real Lives. Real Pakistan.", 0.4, 1.52, 12.6, 0.95, size=50, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

impact = [
    ("240M+","Pakistanis who can benefit from AI healthcare",ACCENT_GREEN),
    ("100K+","Lady Health Workers served by CHW platform",ACCENT_BLUE),
    ("10M+","Deaf Pakistanis served by PSL translator",TEAL),
    ("0 PKR","Cost to diagnose — fully free, open-source",ACCENT_AMBER),
]
for i,(num,desc,col) in enumerate(impact):
    ix = 0.3 + i*3.27
    ib = add_rrect(sl, ix, 2.65, 3.0, 1.7, CARD_BG, col, 1.8, 0.07)
    add_text(sl, num, ix+0.1, 2.7, 2.8, 0.75, size=32, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, desc, ix+0.1, 3.42, 2.8, 0.82, size=9.5, color=WHITE_70, align=PP_ALIGN.CENTER)

add_text(sl, "Roadmap:", 0.35, 4.55, 2.5, 0.3, size=12, color=ACCENT_PURPLE, bold=True)
roadmap = [
    "PWA for Android/iOS offline install",
    "ECG / X-Ray / Ultrasound AI modules",
    "Expansion to Bangladesh, Afghanistan, Indonesia",
    "Punjab/Sindh e-Sehat API integration",
    "Federated Learning — model improvement without data sharing",
]
add_multiline(sl, ["  →  "+r for r in roadmap], 0.35, 4.88, 5.8, 1.5, size=10, color=WHITE_70)

gbox = add_rrect(sl, 6.55, 4.45, 6.5, 2.05, CARD_BG, ACCENT_GREEN, 1.5, 0.06)
add_text(sl, "Links & Resources", 6.7, 4.52, 6.2, 0.34, size=12, color=ACCENT_GREEN, bold=True)
links = [("GitHub:","github.com/hamza0312615/visiondxmega",ACCENT_BLUE),
         ("Dev Server:","npm run dev  →  localhost:5173",ACCENT_GREEN),
         ("CV Backend:","uvicorn main:app  →  localhost:8000",ACCENT_AMBER),
         ("WhatsApp:","node server.js  →  localhost:3001",EMERALD),]
for i,(label,val,lc) in enumerate(links):
    ly = 4.92 + i*0.37
    add_text(sl, label, 6.7, ly, 1.5, 0.33, size=10, color=WHITE_70, bold=True)
    add_text(sl, val, 8.3, ly, 4.6, 0.33, size=10, color=lc, bold=True)

add_text(sl, '"Technology as a bridge — not a barrier — to healthcare."',
         0.3, 6.6, 12.8, 0.42, size=14, color=ACCENT_GREEN, bold=True, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
output_path = r"f:\ALL projects\VisionDX-Mega\VisionDX_Mega_Hackathon_Presentation.pptx"
prs.save(output_path)
print("DONE! Saved to: " + output_path)
print("Total slides: " + str(len(prs.slides)))
